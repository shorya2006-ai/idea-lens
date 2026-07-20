from importlib.metadata import metadata
import json
import pandas as pd
import plotly.express as px
import streamlit as st
import os
import pdfplumber
import time
data_folder = "pdf_data"
os.makedirs(data_folder, exist_ok=True)
from collections import Counter
from pathlib import Path
from datetime import datetime
from groq_helper import generate_groq_analysis
from docx import Document
from pptx import Presentation
from vector_db import search_similar, add_new_idea, initialize_from_pdfs, load_index
# ---------------- INIT ----------------
initialize_from_pdfs()
st.set_page_config(page_title="AI Idea Detector", layout="wide")
if "page" not in st.session_state:
st.session_state.page = "dashboard"
if "user_type" not in st.session_state:
st.session_state.user_type = None
if "search_history" not in st.session_state:
st.session_state.search_history = []
if "selected_review_file" not in st.session_state:
st.session_state.selected_review_file = None
if "pending_status_action" not in st.session_state:
st.session_state.pending_status_action = None
# ---------------- EMPLOYEE PDF ----------------
def load_employee_ids_from_pdf(pdf_path="employee_database_100.pdf"):
emp_ids = set()
try:
with pdfplumber.open(pdf_path) as pdf:
for page in pdf.pages:
text = page.extract_text()
if text:
for line in text.split("\n"):
parts = line.split()
if parts and parts[0].startswith("EMP"):
emp_ids.add(parts[0])
except:
pass
return emp_ids
if "pending_status_file" not in st.session_state:
st.session_state.pending_status_file = None
# ---------------- ADMIN PDF ----------------
def load_admin_credentials_from_pdf(pdf_path="admin_database_5.pdf"):
admins = {}
try:
with pdfplumber.open(pdf_path) as pdf:
for page in pdf.pages:
text = page.extract_text()
for line in text.split("\n"):
parts = line.split()
if len(parts) >= 4 and parts[0].startswith("ADM"):
admins[parts[0]] = {
"name": parts[1] + " " + parts[2],
"password": parts[3]
}
except:
pass
return admins
# ---------------- TEXT EXTRACTION ----------------
def extract_text_from_file(uploaded_file):
text = ""
try:
filename = uploaded_file.name.lower()
if filename.endswith(".pdf"):
with pdfplumber.open(uploaded_file) as pdf:
for page in pdf.pages:
text += page.extract_text() or ""
elif filename.endswith(".docx"):
doc = Document(uploaded_file)
for para in doc.paragraphs:
text += para.text + "\n"
elif filename.endswith(".pptx"):
prs = Presentation(uploaded_file)
for slide in prs.slides:
for shape in slide.shapes:
if hasattr(shape, "text"):
text += shape.text + "\n"
else:
text = uploaded_file.read().decode(errors="ignore")
except:
text = ""
return text
# ---------------- AI EXPLANATION ----------------
def generate_ai_explanation(input_text, matched_text, score):
try:
if not matched_text:
return "No matched content available to generate explanation."
input_words = set(input_text.lower().split())
matched_words = set(matched_text.lower().split())
common_words = list(input_words.intersection(matched_words))
if score > 0.75:
level = "highly similar"
elif score > 0.5:
level = "moderately similar"
else:
level = "mostly different"
explanation = f"The ideas are {level}."
if common_words:
explanation += " Common keywords include: " + ", ".join(common_words[:8])
else:
explanation += " There are no strong overlapping keywords but conceptual similarity exists."
return explanation
except:
return "AI explanation could not be generated"
# ---------- CONTRIBUTOR ANALYTICS ----------
CONTRIBUTOR_FILE = "contributions.json"
def load_contributions():
try:
if os.path.exists(CONTRIBUTOR_FILE):
with open(CONTRIBUTOR_FILE, "r") as f:
return json.load(f)
except:
pass
return {}
def save_contributions(data):
try:
with open(CONTRIBUTOR_FILE, "w") as f:
json.dump(data, f, indent=4)
except:
pass
def increment_contribution(emp_id):
try:
data = load_contributions()
data[emp_id] = data.get(emp_id, 0) + 1
save_contributions(data)
except:
pass
def get_top_contributors():
_, metadata = load_index()
contributors = []
for item in metadata:
source = item.get("source", "")
if "(" in source and ")" in source:
employee = (
source.split("(")[-1]
.replace(")", "")
.strip()
)
contributors.append(employee)
counts = Counter(contributors)
return sorted(
counts.items(),
key=lambda x: x[1],
reverse=True
SEARCH_HISTORY_FILE = "search_history.json"
def load_search_history():
if not os.path.exists(SEARCH_HISTORY_FILE):
return {}
try:
with open(SEARCH_HISTORY_FILE, "r") as f:
return json.load(f)
except:
return {}
#----------------IDEA STATUS--------------
IDEA_STATUS_FILE = "idea_status.json"
def load_idea_status():
if not os.path.exists(IDEA_STATUS_FILE):
return {}
try:
with open(IDEA_STATUS_FILE, "r") as f:
return json.load(f)
except:
return {}
def save_idea_status(data):
with open(IDEA_STATUS_FILE, "w") as f:
json.dump(data, f, indent=4)
def add_idea_status(file_name, employee):
data = load_idea_status()
data[file_name] = {
"employee": employee,
"status": "Pending",
"viewed": False,
"employee_notified": True
}
save_idea_status(data)
def update_idea_status(file_name, status):
data = load_idea_status()
if file_name in data:
data[file_name]["status"] = status
data[file_name]["viewed"] = True
data[file_name]["employee_notified"] = False
save_idea_status(data)
def mark_notification_as_read(file_name):
data = load_idea_status()
if file_name in data:
data[file_name]["employee_notified"] = True
save_idea_status(data)
#----------Search History--------
SEARCH_HISTORY_FILE = "search_history.json"
def load_search_history():
if not os.path.exists(SEARCH_HISTORY_FILE):
return {}
try:
with open(SEARCH_HISTORY_FILE, "r") as f:
return json.load(f)
except:
return {}
def save_search_history(data):
with open(SEARCH_HISTORY_FILE, "w") as f:
json.dump(data, f, indent=4)
def add_search_history(emp_id, search_text):
data = load_search_history()
if emp_id not in data:
data[emp_id] = []
data[emp_id].append(search_text)
save_search_history(data)
# ---------------- DASHBOARD ----------------
if st.session_state.page == "dashboard":
st.markdown("<h1 style='text-align:center;'>AI Idea Duplicate Detector</h1>", unsafe_allow_html=True)
col1, col2 = st.columns(2)
with col1:
if st.button("Admin Login", use_container_width=True):
st.session_state.page = "admin_login"
with col2:
if st.button("Employee Login", use_container_width=True):
st.session_state.page = "employee_login"
# ---------------- ADMIN LOGIN ----------------
elif st.session_state.page == "admin_login":
st.title("Admin Login")
admin_id = st.text_input("Admin ID")
password = st.text_input("Password", type="password")
if st.button("Login"):
admins = load_admin_credentials_from_pdf()
if admin_id in admins and password == admins[admin_id]["password"]:
st.session_state.user_type = "admin"
st.session_state.admin_name = admins[admin_id]["name"]
st.session_state.page = "main"
else:
st.error("Invalid credentials")
# ---------------- EMPLOYEE LOGIN ----------------
elif st.session_state.page == "employee_login":
st.title("Employee Login")
emp_id = st.text_input("Employee ID")
if st.button("Login"):
valid_ids = load_employee_ids_from_pdf()
if emp_id in valid_ids:
st.session_state.user_type = "employee"
st.session_state.emp_id = emp_id
st.session_state.page = "main"
else:
st.error("Invalid Employee ID")
# ---------------- MAIN ----------------
elif st.session_state.page == "main":
# Top bar (Logout right)
col1, col2 = st.columns([10, 1])
with col1:
st.title("AI Idea Duplicate Detector")
with col2:
if st.button("Logout"):
st.session_state.page = "dashboard"
st.session_state.user_type = None
st.rerun()
